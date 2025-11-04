from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def get_template_themes():
    """
    Returns a dictionary of 18 professional presentation themes,
    each with its own unique assigned design style.
    """
    return {
        # Original 6 Designs
        "Modern Corporate": {"background": "#FFFFFF", "font_color": "#000000", "accent_color": "#00529B", "design_style": "sidebar"},
        "Creative Blue": {"background": "#EBF5FF", "font_color": "#003366", "accent_color": "#00AEEF", "design_style": "corner_accent"},
        "Elegant Dark": {"background": "#2C3E50", "font_color": "#FFFFFF", "accent_color": "#F39C12", "design_style": "header_lines"},
        "Professional Green": {"background": "#F0FFF0", "font_color": "#006400", "accent_color": "#3CB371", "design_style": "gradient_bar"},
        "Minimalist Grey": {"background": "#F5F5F5", "font_color": "#333333", "accent_color": "#777777", "design_style": "geometric_shapes"},
        "Bold Red": {"background": "#FFEEEE", "font_color": "#8B0000", "accent_color": "#FF4136", "design_style": "border_frame"},
        
        # New Unique Designs
        "Sunny Yellow": {"background": "#FFFFF0", "font_color": "#3D2B1F", "accent_color": "#FFC300", "design_style": "title_underline"},
        "Deep Purple": {"background": "#F3F0FF", "font_color": "#4B0082", "accent_color": "#9370DB", "design_style": "pillar_bars"},
        "Tech Slate": {"background": "#34495E", "font_color": "#ECF0F1", "accent_color": "#1ABC9C", "design_style": "diagonal_accent"},
        "Warm Orange": {"background": "#FFF5E1", "font_color": "#8B4513", "accent_color": "#FF851B", "design_style": "bottom_left_block"},
        "Classic Navy": {"background": "#E6E6FA", "font_color": "#000080", "accent_color": "#4169E1", "design_style": "top_right_triangle"},
        "Vibrant Magenta": {"background": "#FFF0F5", "font_color": "#8B008B", "accent_color": "#FF00FF", "design_style": "offset_blocks"},
        "Oceanic Teal": {"background": "#F0FFFF", "font_color": "#008080", "accent_color": "#20B2AA", "design_style": "faded_circle"},
        "Earthy Tones": {"background": "#FAF0E6", "font_color": "#5D4037", "accent_color": "#A0522D", "design_style": "chevron_bottom"},
        "Sunset Orange": {"background": "#FFF8DC", "font_color": "#8B4513", "accent_color": "#FF6347", "design_style": "top_bottom_thin_bars"},
        "Minty Fresh": {"background": "#F5FFFA", "font_color": "#2E8B57", "accent_color": "#66CDAA", "design_style": "left_edge_wave"},
        "Royal Gold": {"background": "#1A1A1A", "font_color": "#EAEAEA", "accent_color": "#FFD700", "design_style": "dotted_background"},
        "Monochrome Noir": {"background": "#FFFFFF", "font_color": "#000000", "accent_color": "#808080", "design_style": "cross_lines"},
    }

def apply_theme(prs, theme_name):
    """Applies a dynamic theme to the presentation."""
    themes = get_template_themes()
    theme = themes.get(theme_name, themes["Modern Corporate"])

    # Slide Master
    slide_master = prs.slide_master
    fill = slide_master.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(theme["background"][1:])

    # Title Slide Layout
    title_layout = prs.slide_layouts[0]
    title_placeholder = title_layout.shapes.title
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["font_color"][1:])
    
    # Content Slide Layout
    content_layout = prs.slide_layouts[1]
    title_placeholder = content_layout.shapes.title
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["accent_color"][1:])
    body_placeholder = content_layout.shapes.placeholders[1]
    body_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["font_color"][1:])

    return prs