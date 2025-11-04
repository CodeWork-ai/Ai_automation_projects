from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import traceback
from dynamic_template_generator import get_template_themes

# --- HELPER FUNCTIONS ---
def set_slide_background(slide, theme):
    """Sets the background color for a single slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(theme["background"][1:])

def _set_fill(shape, theme, key="accent_color", transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme[key][1:])
    shape.fill.transparency = transparency
    shape.line.fill.background()

# --- DESIGN FUNCTION LIBRARY (18 UNIQUE DESIGNS) ---
def add_sidebar(slide, theme):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(0.15), Inches(8.6))
    _set_fill(shape, theme)

def add_corner_accent(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(14), Inches(0.3), Inches(1.7), Inches(0.15)), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(15.55), Inches(0.3), Inches(0.15), Inches(1.2)), theme)

def add_header_lines(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.0), Inches(14), Inches(0.04)), theme)

def add_gradient_bar(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.8), Inches(8), Inches(0.2)), theme, "accent_color")
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(8.8), Inches(8), Inches(0.2)), theme, "font_color")

def add_geometric_shapes(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4)), theme, transparency=0.85)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(14), Inches(7), Inches(2.5), Inches(2.5)), theme, transparency=0.85)

def add_border_frame(slide, theme):
    t = Inches(0.08)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), t), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(9) - t, Inches(16), t), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, t, Inches(9)), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(16) - t, 0, t, Inches(9)), theme)

def add_title_underline(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(5), Inches(0.1)), theme)

def add_pillar_bars(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.2), Inches(8)), theme, transparency=0.5)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(15.3), Inches(0.5), Inches(0.2), Inches(8)), theme, transparency=0.5)

def add_diagonal_accent(slide, theme):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(7), Inches(5), Inches(2.5))
    shape.rotation = -15
    _set_fill(shape, theme)

def add_bottom_left_block(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7), Inches(3), Inches(2)), theme)

def add_top_right_triangle(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(13), 0, Inches(3), Inches(3)), theme)

def add_offset_blocks(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(14), Inches(6), Inches(2), Inches(2)), theme, transparency=0.5)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.5), Inches(6.5), Inches(2), Inches(2)), theme, transparency=0.7)

def add_faded_circle(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-3), Inches(8), Inches(8)), theme, transparency=0.8)

def add_chevron_bottom(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6), Inches(8.5), Inches(4), Inches(0.5)), theme)

def add_top_bottom_thin_bars(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.1), Inches(16), Inches(0.05)), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(8.85), Inches(16), Inches(0.05)), theme)

def add_left_edge_wave(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(0), Inches(1.5), Inches(3)), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(3), Inches(1.5), Inches(3)), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(6), Inches(1.5), Inches(3)), theme)

def add_dotted_background(slide, theme):
    for x in range(0, 16, 2):
        for y in range(0, 9, 2):
            _set_fill(slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.1), Inches(0.1)), theme, transparency=0.75)

def add_cross_lines(slide, theme):
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.45), Inches(14), Inches(0.03)), theme)
    _set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.95), Inches(1), Inches(0.03), Inches(7)), theme, transparency=0.5)

# --- DESIGN DISPATCHER ---
DESIGN_MAP = {
    "sidebar": add_sidebar, "corner_accent": add_corner_accent, "header_lines": add_header_lines,
    "gradient_bar": add_gradient_bar, "geometric_shapes": add_geometric_shapes, "border_frame": add_border_frame,
    "title_underline": add_title_underline, "pillar_bars": add_pillar_bars, "diagonal_accent": add_diagonal_accent,
    "bottom_left_block": add_bottom_left_block, "top_right_triangle": add_top_right_triangle, "offset_blocks": add_offset_blocks,
    "faded_circle": add_faded_circle, "chevron_bottom": add_chevron_bottom, "top_bottom_thin_bars": add_top_bottom_thin_bars,
    "left_edge_wave": add_left_edge_wave, "dotted_background": add_dotted_background, "cross_lines": add_cross_lines,
}

def create_presentation(structure, theme_name, output_path):
    try:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(16), Inches(9)
        themes = get_template_themes()
        theme = themes.get(theme_name, themes["Modern Corporate"])
        design_style = theme.get("design_style")
        design_func = DESIGN_MAP.get(design_style)

        for slide_data in structure:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_background(slide, theme)
            
            if design_func:
                design_func(slide, theme)

            # --- CONTENT PLACEMENT (remains the same) ---
            slide_type = slide_data.get('slide_type', 'content_slide')
            content = slide_data.get('content', {})
            
            if slide_type == 'title_slide':
                title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2)).text_frame
                p = title_box.paragraphs[0]; p.text = content.get('title', 'Presentation Title')
                p.font.bold = True; p.font.size = Pt(60); p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = RGBColor.from_string(theme["font_color"][1:])
                
                subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(1)).text_frame
                p = subtitle_box.paragraphs[0]; p.text = content.get('subtitle', '')
                p.font.size = Pt(32); p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = RGBColor.from_string(theme["accent_color"][1:])

            elif slide_type == 'thank_you_slide':
                title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2)).text_frame
                p = title_box.paragraphs[0]; p.text = content.get('title', 'Thank You')
                p.font.bold = True; p.font.size = Pt(66); p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = RGBColor.from_string(theme["font_color"][1:])
                
                contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(1.5)).text_frame
                p1 = contact_box.add_paragraph(); p1.text = "+91 81481 01923"; p1.font.size = Pt(32); p1.alignment = PP_ALIGN.CENTER
                p1.font.color.rgb = RGBColor.from_string(theme["accent_color"][1:])
                p2 = contact_box.add_paragraph(); p2.text = "sales@codework.ai"; p2.font.size = Pt(32); p2.alignment = PP_ALIGN.CENTER
                p2.font.color.rgb = RGBColor.from_string(theme["accent_color"][1:])

            else: # content_slide
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5)).text_frame
                p = title_box.paragraphs[0]; p.text = content.get('title', 'Slide Title')
                p.font.bold = True; p.font.size = Pt(44)
                p.font.color.rgb = RGBColor.from_string(theme["accent_color"][1:])
                
                # Special handling for designs that are tied to the title position
                if design_style == "header_lines" or design_style == "title_underline":
                    design_func(slide, theme)

                body_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(13.5), Inches(6)).text_frame
                for point in content.get('bullet_points', []):
                    p = body_box.add_paragraph(); p.text = f"•  {str(point)}"
                    p.font.size = Pt(28); p.level = 0; p.space_after = Pt(12)
                    p.font.color.rgb = RGBColor.from_string(theme["font_color"][1:])

        prs.save(output_path)
        return output_path
    except Exception as e:
        traceback.print_exc()
        return None